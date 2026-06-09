"""Fallback parsers for file types not supported by MinerU API or when MinerU API fails."""

import os


def parse_text(file_path):
    """Parse plain text files."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(file_path, "r", encoding="latin-1") as f:
            return f.read()


def parse_csv(file_path):
    """Parse CSV/TSV files."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(file_path, "r", encoding="latin-1") as f:
            return f.read()


def parse_docx(file_path):
    """Parse Word documents using python-docx."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError(
            "python-docx is required to parse .docx files. "
            "Install it with: pip install python-docx"
        )
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def parse_pptx(file_path):
    """Parse PowerPoint files using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx is required to parse .pptx files. "
            "Install it with: pip install python-pptx"
        )
    prs = Presentation(file_path)
    slides_text = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text:
                        texts.append(para.text)
        slides_text.append("\n".join(texts))
    return "\n---\n".join(slides_text)


def parse_image(file_path):
    """Parse image files using Pillow."""
    filename = os.path.basename(file_path)
    try:
        from PIL import Image
        with Image.open(file_path) as img:
            return (
                f"Image: {filename}\n"
                f"Format: {img.format}\n"
                f"Size: {img.width}x{img.height}\n"
                f"Mode: {img.mode}"
            )
    except ImportError:
        size = os.path.getsize(file_path)
        return f"Image: {filename}\nSize: {size} bytes"


def parse_excel(file_path):
    """Parse Excel files using openpyxl."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"Sheet: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                parts.append("\t".join(str(c) if c is not None else "" for c in row))
        wb.close()
        return "\n".join(parts)
    except ImportError:
        size = os.path.getsize(file_path)
        filename = os.path.basename(file_path)
        return f"Excel: {filename}\nSize: {size} bytes"


_TEXT_EXTS = {
    ".txt", ".md", ".log", ".py", ".js", ".json", ".xml", ".html", ".css",
    ".yaml", ".yml", ".ini", ".cfg", ".conf", ".sh", ".bat", ".sql", ".r",
    ".java", ".c", ".cpp", ".h", ".hpp", ".go", ".rs", ".rb", ".php",
    ".swift", ".kt", ".ts", ".tsx", ".jsx", ".vue", ".svelte",
}

_CSV_EXTS = {".csv", ".tsv"}

_DOCX_EXTS = {".docx"}

_PPTX_EXTS = {".pptx"}

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}

_EXCEL_EXTS = {".xlsx", ".xls"}

_EXT_MAP = {}
for ext in _TEXT_EXTS:
    _EXT_MAP[ext] = parse_text
for ext in _CSV_EXTS:
    _EXT_MAP[ext] = parse_csv
for ext in _DOCX_EXTS:
    _EXT_MAP[ext] = parse_docx
for ext in _PPTX_EXTS:
    _EXT_MAP[ext] = parse_pptx
for ext in _IMAGE_EXTS:
    _EXT_MAP[ext] = parse_image
for ext in _EXCEL_EXTS:
    _EXT_MAP[ext] = parse_excel


def get_fallback_parser(file_path):
    """Return the appropriate fallback parser function for a given file path based on its extension."""
    _, ext = os.path.splitext(file_path)
    return _EXT_MAP.get(ext.lower())
